"""Document ingestion pipeline: extract text from Literature/ files, chunk, and store in ChromaDB."""

import os
import re
import hashlib
from pathlib import Path

import chromadb
import pdfplumber
from pptx import Presentation
from docx import Document as DocxDocument
from openpyxl import load_workbook
import tiktoken

TOPIC_MAP = {
    "1": "საბიუჯეტო კანონმდებლობა",
    "2": "ბუღალტრული აღრიცხვა საჯარო სექტორში",
    "3": "ფინანსური აუდიტი",
    "4": "შესაბამისობის აუდიტი",
    "5": "ეფექტიანობის აუდიტი",
    "general": "ზოგადი მასალა",
}

CHUNK_SIZE = 800
CHUNK_OVERLAP = 200

enc = tiktoken.encoding_for_model("gpt-4o")


def token_len(text: str) -> int:
    return len(enc.encode(text))


def chunk_text(text: str, max_tokens: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    paragraphs = [p.strip() for p in re.split(r"\n{2,}", text) if p.strip()]
    chunks: list[str] = []
    current: list[str] = []
    current_len = 0

    for para in paragraphs:
        para_len = token_len(para)
        if current_len + para_len > max_tokens and current:
            chunks.append("\n\n".join(current))
            overlap_text = "\n\n".join(current)
            overlap_tokens = token_len(overlap_text)
            while current and overlap_tokens > overlap:
                current.pop(0)
                overlap_text = "\n\n".join(current)
                overlap_tokens = token_len(overlap_text)
            current_len = overlap_tokens
        current.append(para)
        current_len += para_len

    if current:
        chunks.append("\n\n".join(current))

    return chunks


def detect_topic(path: str) -> tuple[str, str]:
    parts = path.replace("\\", "/").split("/")
    for part in parts:
        for key, name in TOPIC_MAP.items():
            if key == "general":
                continue
            if part.startswith(f"{key}."):
                return key, name
    for part in parts:
        if "General" in part:
            return "general", TOPIC_MAP["general"]
        if "Testchats" in part or "testchat" in part.lower():
            return "general", "ტესტის კითხვები"
    return "general", TOPIC_MAP["general"]


def extract_pdf(filepath: str) -> str:
    texts = []
    with pdfplumber.open(filepath) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                texts.append(t)
    return "\n\n".join(texts)


def extract_pptx(filepath: str) -> str:
    prs = Presentation(filepath)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
    return "\n\n".join(texts)


def extract_docx(filepath: str) -> str:
    doc = DocxDocument(filepath)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_xlsx(filepath: str) -> str:
    wb = load_workbook(filepath, data_only=True)
    texts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            vals = [str(c) for c in row if c is not None]
            if vals:
                texts.append(" | ".join(vals))
    return "\n\n".join(texts)


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
}


def ingest(literature_path: str, db_path: str = "./chroma_db") -> dict:
    client = chromadb.PersistentClient(path=db_path)
    try:
        client.delete_collection("audit_docs")
    except Exception:
        pass
    collection = client.get_or_create_collection(
        name="audit_docs",
        metadata={"hnsw:space": "cosine"},
    )

    stats = {"files_processed": 0, "chunks_created": 0, "errors": []}

    for root, _dirs, files in os.walk(literature_path):
        for fname in sorted(files):
            ext = Path(fname).suffix.lower()
            if ext not in EXTRACTORS:
                continue
            filepath = os.path.join(root, fname)
            topic_id, topic_name = detect_topic(filepath)

            try:
                text = EXTRACTORS[ext](filepath)
                if not text.strip():
                    stats["errors"].append(f"Empty: {fname}")
                    continue

                chunks = chunk_text(text)
                for i, chunk in enumerate(chunks):
                    doc_id = hashlib.md5(f"{filepath}:{i}".encode()).hexdigest()
                    collection.add(
                        ids=[doc_id],
                        documents=[chunk],
                        metadatas=[{
                            "source": fname,
                            "topic_id": topic_id,
                            "topic_name": topic_name,
                            "chunk_index": i,
                        }],
                    )
                stats["files_processed"] += 1
                stats["chunks_created"] += len(chunks)
                print(f"  ✓ {fname} → {len(chunks)} chunks [{topic_name}]")

            except Exception as e:
                stats["errors"].append(f"{fname}: {e}")
                print(f"  ✗ {fname}: {e}")

    return stats


if __name__ == "__main__":
    import sys
    lit_path = sys.argv[1] if len(sys.argv) > 1 else "../../Literature"
    print(f"Ingesting from: {lit_path}")
    result = ingest(lit_path)
    print(f"\nDone: {result['files_processed']} files, {result['chunks_created']} chunks")
    if result["errors"]:
        print(f"Errors: {len(result['errors'])}")
        for e in result["errors"]:
            print(f"  - {e}")
